"""
Mail2Task_Presentation_v2_recommendation.pptx 에
Apriori 알고리즘 상세 슬라이드를 추가합니다.
슬라이드 6 (동작 구조) 바로 뒤에 삽입됩니다.

실행: python -m scripts.add_apriori_slide
"""
from __future__ import annotations

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "Mail2Task_Presentation_v2_recommendation.pptx"

EMU_PER_INCH = 914_400


# ── helpers ──────────────────────────────────────────────────────────────

def rgb(h: str) -> RGBColor:
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _emu(inches: float) -> int:
    return int(inches * EMU_PER_INCH)


def _no_border(shape) -> None:
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for child in list(ln):
        ln.remove(child)
    etree.SubElement(ln, qn("a:noFill"))


def rect(slide, x: float, y: float, w: float, h: float, fill_hex: str):
    """Filled rectangle with no visible border."""
    s = slide.shapes.add_shape(1, _emu(x), _emu(y), _emu(w), _emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_hex)
    _no_border(s)
    return s


def rect_text(slide, x, y, w, h, fill_hex, lines, center=False):
    """Filled rectangle with styled text lines."""
    s = rect(slide, x, y, w, h, fill_hex)
    tf = s.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 10))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = rgb(ln.get("color", "FFFFFF"))
    return s


def label(slide, x, y, w, h, lines):
    """Transparent text box with styled lines."""
    tb = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 10))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = rgb(ln.get("color", "000000"))
    return tb


def move_slide(prs, old_index: int, new_index: int) -> None:
    sl = prs.slides._sldIdLst
    els = list(sl)
    el = els[old_index]
    sl.remove(el)
    sl.insert(new_index, el)


# ── main ─────────────────────────────────────────────────────────────────

def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # 레이아웃 placeholder 제거
    for ph in list(slide.placeholders):
        ph.element.getparent().remove(ph.element)

    # ── 헤더 바 ───────────────────────────────────────────────
    rect(slide, 0, 0, 10, 0.82, "1B2A4A")
    rect_text(slide, 0.38, 0.15, 0.62, 0.52, "7C3AED",
              [{"text": "04", "size": 15, "bold": True, "color": "FFFFFF"}],
              center=True)
    label(slide, 1.18, 0.06, 8.55, 0.70, [
        {"text": "PDF 연관 추천 — Apriori 알고리즘 상세",
         "size": 20, "bold": True, "color": "FFFFFF"},
    ])

    # ── 서브타이틀 바 ─────────────────────────────────────────
    rect(slide, 0, 0.82, 10, 0.35, "EBF5F8")
    rect(slide, 0, 0.82, 0.05, 0.35, "7C3AED")
    label(slide, 0.22, 0.82, 9.4, 0.35, [
        {"text": "장바구니 분석(Market Basket) 기법으로 PDF 간 연관도를 자동 계산하는 방식",
         "size": 11, "color": "64748B"},
    ])

    # ── 페이지 카운터 ──────────────────────────────────────────
    rect_text(slide, 8.0, 5.33, 1.7, 0.22, "F5F7FA",
              [{"text": "4  /  6", "size": 10, "color": "64748B"}])

    # ── STEP 1: 토큰 추출 (blue) ──────────────────────────────
    rect(slide, 0.45, 1.22, 2.80, 1.58, "065A82")
    rect_text(slide, 0.55, 1.29, 0.30, 0.30, "FFFFFF",
              [{"text": "1", "size": 9, "bold": True, "color": "065A82"}], center=True)
    label(slide, 0.93, 1.27, 2.18, 0.30, [
        {"text": "STEP 1  ·  토큰 추출", "size": 9, "bold": True, "color": "7DD3FC"},
    ])
    label(slide, 0.93, 1.57, 2.18, 0.28, [
        {"text": "📄 각 PDF = 하나의 트랜잭션", "size": 10, "bold": True, "color": "FFFFFF"},
    ])
    label(slide, 0.93, 1.85, 2.18, 0.87, [
        {"text": '정규식으로 키워드 상위 20개 추출', "size": 8.5, "color": "BAE6FD"},
        {"text": '예) ["ci", "cd", "pipeline", "build"]', "size": 8.5, "color": "BAE6FD"},
    ])

    # 화살표 1
    rect_text(slide, 3.28, 1.87, 0.18, 0.26, "0D9488",
              [{"text": "▶", "size": 10, "bold": True, "color": "FFFFFF"}], center=True)

    # ── STEP 2: Apriori 연관 규칙 (purple) ───────────────────
    rect(slide, 3.50, 1.22, 2.80, 1.58, "7C3AED")
    rect_text(slide, 3.60, 1.29, 0.30, 0.30, "FFFFFF",
              [{"text": "2", "size": 9, "bold": True, "color": "7C3AED"}], center=True)
    label(slide, 3.97, 1.27, 2.18, 0.30, [
        {"text": "STEP 2  ·  연관 규칙 생성", "size": 9, "bold": True, "color": "DDD6FE"},
    ])
    label(slide, 3.97, 1.57, 2.18, 0.28, [
        {"text": "🔗 Apriori 알고리즘 실행", "size": 10, "bold": True, "color": "FFFFFF"},
    ])
    label(slide, 3.97, 1.85, 2.18, 0.87, [
        {"text": "{pipeline, docker} → 규칙 생성", "size": 8.5, "color": "EDE9FE"},
        {"text": "support / confidence / lift 계산", "size": 8.5, "color": "EDE9FE"},
    ])

    # 화살표 2
    rect_text(slide, 6.33, 1.87, 0.18, 0.26, "0D9488",
              [{"text": "▶", "size": 10, "bold": True, "color": "FFFFFF"}], center=True)

    # ── STEP 3: 최종 점수 (teal) ──────────────────────────────
    rect(slide, 6.55, 1.22, 3.00, 1.58, "0D9488")
    rect_text(slide, 6.65, 1.29, 0.30, 0.30, "FFFFFF",
              [{"text": "3", "size": 9, "bold": True, "color": "0D9488"}], center=True)
    label(slide, 7.02, 1.27, 2.38, 0.30, [
        {"text": "STEP 3  ·  최종 점수 산출", "size": 9, "bold": True, "color": "A7F3D0"},
    ])
    label(slide, 7.02, 1.57, 2.38, 0.28, [
        {"text": "📊 두 지표 결합 후 추천", "size": 10, "bold": True, "color": "FFFFFF"},
    ])
    label(slide, 7.02, 1.85, 2.38, 0.87, [
        {"text": "score = Apriori×2.5 + Jaccard", "size": 8.5, "color": "D1FAE5"},
        {"text": "Jaccard = |A∩B| / |A∪B|",        "size": 8.5, "color": "D1FAE5"},
    ])

    # ── 다크 패널 (하단 상세 설명) ────────────────────────────
    rect(slide, 0.45, 2.97, 9.10, 2.18, "1B2A4A")
    rect(slide, 0.45, 2.97, 0.05, 2.18, "7C3AED")
    # 컬럼 구분선
    rect(slide, 3.55, 3.07, 0.04, 1.97, "374151")
    rect(slide, 6.55, 3.07, 0.04, 1.97, "374151")

    # 컬럼 1: 3가지 연관 지표
    label(slide, 0.62, 3.03, 2.80, 0.36, [
        {"text": "📐  3가지 연관 지표", "size": 11, "bold": True, "color": "C4B5FD"},
    ])
    label(slide, 0.62, 3.42, 2.80, 1.62, [
        {"text": "Support",    "size": 9.5, "bold": True,  "color": "7DD3FC"},
        {"text": "  전체 PDF 중 키워드 조합 포함 비율", "size": 8.5, "color": "D1D5DB"},
        {"text": "Confidence", "size": 9.5, "bold": True,  "color": "7DD3FC", "space_before": 3},
        {"text": "  소스→후보 PDF 조건부 확률",        "size": 8.5, "color": "D1D5DB"},
        {"text": "Lift",       "size": 9.5, "bold": True,  "color": "7DD3FC", "space_before": 3},
        {"text": "  1.0 이상이면 의미 있는 연관",      "size": 8.5, "color": "D1D5DB"},
    ])

    # 컬럼 2: 동작 예시
    label(slide, 3.65, 3.03, 2.78, 0.36, [
        {"text": "💡  동작 예시 (PDF 3개 기준)", "size": 11, "bold": True, "color": "6EE7B7"},
    ])
    label(slide, 3.65, 3.42, 2.78, 1.62, [
        {"text": "PDF_A: [ci, cd, pipeline, build]",  "size": 8.5, "color": "D1D5DB"},
        {"text": "PDF_B: [pipeline, docker, deploy]", "size": 8.5, "color": "D1D5DB"},
        {"text": "PDF_C: [api, rest, backend]",       "size": 8.5, "color": "D1D5DB"},
        {"text": "",                                   "size": 4},
        {"text": "→ {pipeline} support = 0.67",       "size": 8.5, "bold": True, "color": "6EE7B7"},
        {"text": "→ {pipeline, build} conf = 0.50",   "size": 8.5, "bold": True, "color": "6EE7B7"},
    ])

    # 컬럼 3: 최종 점수 공식
    label(slide, 6.65, 3.03, 2.78, 0.36, [
        {"text": "🎯  최종 점수 공식", "size": 11, "bold": True, "color": "FDE68A"},
    ])
    label(slide, 6.65, 3.42, 2.78, 1.62, [
        {"text": "규칙 있을 때:",              "size": 9,   "bold": True,  "color": "FDE68A"},
        {"text": "  (규칙합 × 2.5) + Jaccard", "size": 8.5, "color": "D1D5DB"},
        {"text": "규칙 없을 때 (fallback):",   "size": 9,   "bold": True,  "color": "FDE68A", "space_before": 4},
        {"text": "  Jaccard × 1.5",            "size": 8.5, "color": "D1D5DB"},
        {"text": "",                            "size": 4},
        {"text": "→ 상위 5개 PDF 추천 반환",   "size": 8.5, "bold": True,  "color": "86EFAC"},
    ])

    return slide


def main():
    prs = Presentation(str(PPTX_PATH))
    print(f"기존 슬라이드 수: {len(prs.slides)}")

    build_slide(prs)

    # 동작 구조(인덱스 5) 바로 다음인 인덱스 6으로 이동
    move_slide(prs, len(prs.slides) - 1, 6)

    prs.save(str(PPTX_PATH))
    print(f"완료: {PPTX_PATH.name}  (총 {len(prs.slides)}장)")


if __name__ == "__main__":
    main()
