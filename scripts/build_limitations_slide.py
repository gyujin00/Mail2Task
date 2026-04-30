"""
5조_발표자료.pptx — 한계 및 개선사항 슬라이드 완성 스크립트

슬라이드 15 (인덱스 14) 를 인덱스로 직접 지정하여 내용을 교체합니다.
슬라이드 3 (인덱스 2) 의 원래 '프로젝트 배경' 내용도 복원합니다.

실행: python -m scripts.build_limitations_slide
"""
from __future__ import annotations

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "5조_발표자료.pptx"

EMU = 914_400


# ── 공통 헬퍼 ──────────────────────────────────────────────────────────────

def rgb(h: str) -> RGBColor:
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def e(inches: float) -> int:
    return int(inches * EMU)


def _no_border(shape) -> None:
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for c in list(ln):
        ln.remove(c)
    etree.SubElement(ln, qn("a:noFill"))


def rect(slide, x, y, w, h, fill_hex):
    s = slide.shapes.add_shape(1, e(x), e(y), e(w), e(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_hex)
    _no_border(s)
    return s


def rect_text(slide, x, y, w, h, fill_hex, lines, center=False):
    s = rect(slide, x, y, w, h, fill_hex)
    tf = s.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 10))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = rgb(ln.get("color", "FFFFFF"))
    return s


def label(slide, x, y, w, h, lines):
    tb = slide.shapes.add_textbox(e(x), e(y), e(w), e(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 10))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = rgb(ln.get("color", "000000"))
    return tb


def clear_slide(slide) -> None:
    sp_tree = slide.shapes._spTree
    tags = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}
    to_remove = [c for c in sp_tree if c.tag.split("}")[-1] in tags]
    for el in to_remove:
        sp_tree.remove(el)


# ── 슬라이드 15: 한계 및 개선사항 ──────────────────────────────────────────

LIMIT_ITEMS = [
    {
        "num":   "①",
        "title": "메일 처리 범위 한정",
        "line1": "Gmail IMAP 전용 · 최근 50개 메일만 처리",
        "line2": "Outlook · 네이버 등 타 서비스 미지원",
        "c_bar": "991B1B",
        "c_bg":  "FFF1F2",
        "c_ttl": "9F1239",
    },
    {
        "num":   "②",
        "title": "구조화 양식 의존",
        "line1": "'과업명:' 필드 기반 규칙 추출에 한정",
        "line2": "비정형 자연어 메일의 업무 추출 미흡",
        "c_bar": "9A3412",
        "c_bg":  "FFF7ED",
        "c_ttl": "9A3412",
    },
    {
        "num":   "③",
        "title": "마감일 파싱 한계",
        "line1": "정규식 패턴 매칭 방식으로 처리",
        "line2": "'분기말', 'Q3까지' 등 추상 표현 미인식",
        "c_bar": "92400E",
        "c_bg":  "FFFBEB",
        "c_ttl": "92400E",
    },
    {
        "num":   "④",
        "title": "연관 추천 데이터 의존성",
        "line1": "Apriori — PDF 수 부족 시 규칙 미생성",
        "line2": "소량 데이터에서 Jaccard 단독 fallback",
        "c_bar": "6D28D9",
        "c_bg":  "F5F3FF",
        "c_ttl": "5B21B6",
    },
    {
        "num":   "⑤",
        # 사용자 작성 내용 반영 — ML/파인튜닝 시도했으나 GPU·데이터·시간 제약
        "title": "로컬 AI 모델 학습 한계",
        "line1": "규칙 기반 외 ML/파인튜닝 모델 도입 시도",
        "line2": "GPU 성능 부족 · 데이터 부족 · 2일 제한으로 미완",
        "c_bar": "1E3A8A",
        "c_bg":  "EFF6FF",
        "c_ttl": "1E40AF",
    },
]

IMPROVE_ITEMS = [
    {
        "title": "다양한 메일 서비스 지원",
        "line1": "Outlook · 네이버 메일 IMAP 확장",
        "line2": "OAuth 2.0 인증으로 보안 강화",
        "c_bar": "166534",
        "c_bg":  "F0FDF4",
        "c_ttl": "15803D",
    },
    {
        "title": "LLM 기반 비정형 업무 추출",
        "line1": "Claude/GPT로 자연어 메일에서 업무 추출",
        "line2": "현재 rule-based 추출의 보조 레이어로 추가",
        "c_bar": "065F46",
        "c_bg":  "ECFDF5",
        "c_ttl": "047857",
    },
    {
        "title": "마감일 이해 고도화",
        "line1": "LLM으로 '연말까지', 'Q3 전' 등 해석",
        "line2": "캘린더 연동 · 공휴일 자동 반영",
        "c_bar": "065A82",
        "c_bg":  "F0F9FF",
        "c_ttl": "0369A1",
    },
    {
        "title": "경량 로컬 모델 적용",
        "line1": "데이터 축적 후 경량 fine-tuned 모델 도입",
        "line2": "임베딩(BERT) 기반 유사도로 추천 품질 향상",
        "c_bar": "0D9488",
        "c_bg":  "F0FDFA",
        "c_ttl": "0F766E",
    },
    {
        "title": "실시간 알림 · 피드백 루프",
        "line1": "마감 임박 알림 · Slack / Teams 연동",
        "line2": "사용자 피드백으로 분류·추천 모델 자동 개선",
        "c_bar": "1D4ED8",
        "c_bg":  "EFF6FF",
        "c_ttl": "1E40AF",
    },
]


def build_limitations(slide) -> None:
    # ── 헤더 바 ───────────────────────────────────────────────────
    rect(slide, 0, 0, 10, 0.82, "1B2A4A")
    rect_text(slide, 0.38, 0.15, 0.62, 0.52, "DC2626",
              [{"text": "08", "size": 15, "bold": True, "color": "FFFFFF"}],
              center=True)
    label(slide, 1.18, 0.06, 8.55, 0.70, [
        {"text": "한계 및 개선사항",
         "size": 20, "bold": True, "color": "FFFFFF"},
    ])

    # ── 서브타이틀 바 ─────────────────────────────────────────────
    rect(slide, 0, 0.82, 10, 0.36, "FFF1F2")
    rect(slide, 0, 0.82, 0.05, 0.36, "DC2626")
    label(slide, 0.22, 0.85, 9.4, 0.30, [
        {"text": "현재 시스템의 구현 제약과 향후 발전 방향",
         "size": 11, "color": "64748B"},
    ])

    # ── 섹션 레이블 행 ────────────────────────────────────────────
    SEC_Y = 1.28
    rect(slide, 0.35, SEC_Y, 4.50, 0.36, "FEE2E2")
    rect(slide, 0.35, SEC_Y, 0.05, 0.36, "DC2626")
    label(slide, 0.55, SEC_Y + 0.04, 4.20, 0.28, [
        {"text": "⚠  현재 한계  (구현 범위 기준)",
         "size": 11, "bold": True, "color": "B91C1C"},
    ])
    rect(slide, 5.15, SEC_Y, 4.50, 0.36, "DCFCE7")
    rect(slide, 5.15, SEC_Y, 0.05, 0.36, "16A34A")
    label(slide, 5.35, SEC_Y + 0.04, 4.20, 0.28, [
        {"text": "✓  개선 방향  (향후 적용 계획)",
         "size": 11, "bold": True, "color": "15803D"},
    ])

    # ── 세로 구분선 ───────────────────────────────────────────────
    rect(slide, 4.92, 1.24, 0.06, 4.12, "E2E8F0")

    # ── 카드 공통 치수 ────────────────────────────────────────────
    CARD_H  = 0.68
    GAP     = 0.055
    START_Y = 1.76
    LX      = 0.35
    RX      = 5.15
    CW      = 4.50

    for i, item in enumerate(LIMIT_ITEMS):
        y = START_Y + i * (CARD_H + GAP)
        rect(slide, LX, y, CW, CARD_H, item["c_bg"])
        rect(slide, LX, y, 0.05, CARD_H, item["c_bar"])
        label(slide, LX + 0.16, y + 0.06, CW - 0.20, 0.24, [
            {"text": f"{item['num']}  {item['title']}",
             "size": 10, "bold": True, "color": item["c_ttl"]},
        ])
        label(slide, LX + 0.16, y + 0.32, CW - 0.20, 0.34, [
            {"text": item["line1"], "size": 8.5, "color": "374151"},
            {"text": item["line2"], "size": 8.5, "color": "64748B"},
        ])

    for i, item in enumerate(IMPROVE_ITEMS):
        y = START_Y + i * (CARD_H + GAP)
        rect(slide, RX, y, CW, CARD_H, item["c_bg"])
        rect(slide, RX, y, 0.05, CARD_H, item["c_bar"])
        label(slide, RX + 0.16, y + 0.06, CW - 0.20, 0.24, [
            {"text": f"→  {item['title']}",
             "size": 10, "bold": True, "color": item["c_ttl"]},
        ])
        label(slide, RX + 0.16, y + 0.32, CW - 0.20, 0.34, [
            {"text": item["line1"], "size": 8.5, "color": "374151"},
            {"text": item["line2"], "size": 8.5, "color": "64748B"},
        ])

    # ── 하단 메모 바 ──────────────────────────────────────────────
    FOOTER_Y = START_Y + 5 * (CARD_H + GAP) - GAP + 0.07
    rect(slide, 0.35, FOOTER_Y, 9.30, 0.26, "1E293B")
    label(slide, 0.52, FOOTER_Y + 0.04, 9.00, 0.18, [
        {"text": "현재는 데모·수업용 구현 수준 — 데이터 축적 및 GPU 환경 확보 후 단계적으로 적용 예정",
         "size": 8.5, "color": "94A3B8"},
    ])


# ── 슬라이드 3 복원: 프로젝트 배경 ───────────────────────────────────────────

PAIN_POINTS = [
    {
        "icon": "📥",
        "title": "메일 수동 확인",
        "desc":  "업무 요청·마감 정보를 일일이 메일에서 찾아야 함",
    },
    {
        "icon": "📝",
        "title": "수동 To-do 기록",
        "desc":  "별도 툴(메모장·스프레드시트)에 직접 옮겨 적는 이중 작업",
    },
    {
        "icon": "⏰",
        "title": "마감 누락·지연",
        "desc":  "메일 속에 묻힌 마감일 · 우선순위를 놓치는 사례 빈번",
    },
    {
        "icon": "🗂️",
        "title": "첨부파일 관리 어려움",
        "desc":  "PDF 등 첨부파일과 메일 내용을 연결해 추적하기 어려움",
    },
]


def build_background(slide) -> None:
    # ── 헤더 바 ───────────────────────────────────────────────────
    rect(slide, 0, 0, 10, 0.82, "1B2A4A")
    rect_text(slide, 0.38, 0.15, 0.62, 0.52, "0D9488",
              [{"text": "01", "size": 15, "bold": True, "color": "FFFFFF"}],
              center=True)
    label(slide, 1.18, 0.06, 8.55, 0.70, [
        {"text": "기존의 메일 시스템, 왜 비효율적인가?",
         "size": 20, "bold": True, "color": "FFFFFF"},
    ])

    # ── 서브타이틀 바 ─────────────────────────────────────────────
    rect(slide, 0, 0.82, 10, 0.36, "EBF5F8")
    rect(slide, 0, 0.82, 0.05, 0.36, "0D9488")
    label(slide, 0.22, 0.85, 9.4, 0.30, [
        {"text": "기존의 메일 시스템의 현실적인 한계",
         "size": 11, "color": "64748B"},
    ])

    # ── 페이지 카운터 ──────────────────────────────────────────────
    rect_text(slide, 8.0, 5.33, 1.7, 0.22, "F5F7FA",
              [{"text": "1  /  14", "size": 10, "color": "64748B"}])

    # ── 메인 메시지 ───────────────────────────────────────────────
    rect(slide, 0.45, 1.25, 9.10, 0.60, "EBF5F8")
    rect(slide, 0.45, 1.25, 0.05, 0.60, "0D9488")
    label(slide, 0.65, 1.32, 8.80, 0.46, [
        {"text": "업무 메일을 받아도 To-do로 이어지지 않는다 — 사람이 직접 읽고, 기록하고, 추적해야 한다.",
         "size": 11, "color": "1E293B"},
    ])

    # ── 문제 카드 4개 (2×2 그리드) ───────────────────────────────
    POSITIONS = [
        (0.45, 2.00),   # 좌상
        (5.15, 2.00),   # 우상
        (0.45, 3.55),   # 좌하
        (5.15, 3.55),   # 우하
    ]
    CW, CH = 4.35, 1.38
    COLORS = ["065A82", "7C3AED", "0D9488", "B91C1C"]
    BG_COLORS = ["F0F9FF", "FAF5FF", "F0FDFA", "FFF1F2"]

    for (px, py), item, bar_color, bg_color in zip(
        POSITIONS, PAIN_POINTS, COLORS, BG_COLORS
    ):
        rect(slide, px, py, CW, CH, bg_color)
        rect(slide, px, py, 0.05, CH, bar_color)
        label(slide, px + 0.18, py + 0.12, CW - 0.22, 0.38, [
            {"text": f"{item['icon']}  {item['title']}",
             "size": 13, "bold": True, "color": bar_color},
        ])
        label(slide, px + 0.18, py + 0.55, CW - 0.22, 0.76, [
            {"text": item["desc"], "size": 10.5, "color": "374151"},
        ])

    # ── 하단 전환 메시지 ──────────────────────────────────────────
    rect(slide, 0.45, 5.00, 9.10, 0.28, "1B2A4A")
    rect(slide, 0.45, 5.00, 0.05, 0.28, "0D9488")
    label(slide, 0.65, 5.03, 8.80, 0.22, [
        {"text": "→  Mail2Task: 메일을 자동으로 분석해 업무 목록으로 변환합니다.",
         "size": 10, "bold": True, "color": "6EE7B7"},
    ])


# ── 실행 ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(str(PPTX_PATH))
    total = len(prs.slides)
    print(f"기존 슬라이드 수: {total}")

    # 슬라이드 3 (인덱스 2) 복원
    print("\n[1/2] 슬라이드 3 (프로젝트 배경) 복원 중...")
    slide3 = prs.slides[2]
    clear_slide(slide3)
    build_background(slide3)
    print("     완료")

    # 슬라이드 15 (인덱스 14) 한계 및 개선사항 완성
    print("[2/2] 슬라이드 15 (한계 및 개선사항) 업데이트 중...")
    slide15 = prs.slides[14]
    clear_slide(slide15)
    build_limitations(slide15)
    print("     완료")

    prs.save(str(PPTX_PATH))
    print(f"\n✓ 저장 완료: {PPTX_PATH.name}")
    print("  ※ 슬라이드 3은 원본 내용 복원이 불가해 유사 내용으로 재구성했습니다.")
    print("    내용을 확인 후 필요하면 직접 수정해 주세요.")


if __name__ == "__main__":
    main()
