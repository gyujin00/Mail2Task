"""
슬라이드 7 (인덱스 6) — Apriori 연관분석 슬라이드를
단순화된 원리 설명(지지도/신뢰도/향상도)으로 교체합니다.

실행: python -m scripts.update_apriori_slide
"""
from __future__ import annotations

import sys
sys.stdout.reconfigure(encoding="utf-8")

from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parents[1]
PPTX_PATH = ROOT / "Mail2Task_Presentation_v2_recommendation.pptx"

EMU = 914_400


# ── 공통 헬퍼 ────────────────────────────────────────────────────────────

def rgb(h: str) -> RGBColor:
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def e(inches: float) -> int:
    return int(inches * EMU)


def _no_border(shape) -> None:
    spPr = shape._element.spPr
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for c in list(ln):
        ln.remove(c)
    etree.SubElement(ln, qn("a:noFill"))


def rect(slide, x, y, w, h, fill_hex):
    s = slide.shapes.add_shape(1, e(x), e(y), e(w), e(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill_hex)
    _no_border(s)
    return s


def rect_text(slide, x, y, w, h, fill_hex, lines, center=False):
    s = rect(slide, x, y, w, h, fill_hex)
    tf = s.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 10))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = rgb(ln.get("color", "FFFFFF"))
    return s


def label(slide, x, y, w, h, lines):
    tb = slide.shapes.add_textbox(e(x), e(y), e(w), e(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", PP_ALIGN.LEFT)
        if ln.get("space_before"):
            p.space_before = Pt(ln["space_before"])
        r = p.add_run()
        r.text = ln.get("text", "")
        r.font.size = Pt(ln.get("size", 10))
        r.font.bold = ln.get("bold", False)
        r.font.color.rgb = rgb(ln.get("color", "000000"))
    return tb


def clear_slide(slide) -> None:
    sp_tree = slide.shapes._spTree
    tags = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}
    to_remove = [c for c in sp_tree if c.tag.split("}")[-1] in tags]
    for el in to_remove:
        sp_tree.remove(el)


# ── 슬라이드 빌드 ─────────────────────────────────────────────────────────

def build(slide) -> None:

    # ── 헤더 바 ───────────────────────────────────────────────
    rect(slide, 0, 0, 10, 0.82, "1B2A4A")
    rect_text(slide, 0.38, 0.15, 0.62, 0.52, "7C3AED",
              [{"text": "04", "size": 15, "bold": True}], center=True)
    label(slide, 1.18, 0.06, 8.55, 0.70, [
        {"text": "PDF 연관 추천 — 연관분석(Apriori) 원리",
         "size": 20, "bold": True, "color": "FFFFFF"},
    ])

    # ── 서브타이틀 바 ─────────────────────────────────────────
    rect(slide, 0, 0.82, 10, 0.35, "EBF5F8")
    rect(slide, 0, 0.82, 0.05, 0.35, "7C3AED")
    label(slide, 0.22, 0.82, 9.4, 0.35, [
        {"text": "자주 함께 등장하는 키워드 패턴을 찾아 PDF 간 연관성을 자동으로 계산하는 방법",
         "size": 11, "color": "64748B"},
    ])

    # ── 페이지 카운터 ──────────────────────────────────────────
    rect_text(slide, 8.0, 5.33, 1.7, 0.22, "F5F7FA",
              [{"text": "4  /  6", "size": 10, "color": "64748B"}])

    # ── 인트로 스트립 ──────────────────────────────────────────
    rect(slide, 0.45, 1.22, 9.10, 0.63, "EBF5F8")
    rect(slide, 0.45, 1.22, 0.05, 0.63, "7C3AED")
    label(slide, 0.65, 1.28, 8.70, 0.52, [
        {"text": "각 PDF의 키워드를 '거래 목록'으로 취급합니다. "
                 "여러 PDF에 걸쳐 자주 함께 등장하는 키워드 조합을 찾고,",
         "size": 10.5, "color": "1E293B"},
        {"text": "그 조합이 얼마나 자주(지지도) · 얼마나 믿을 만한지(신뢰도) · "
                 "얼마나 강한 연관인지(향상도) — 세 지표로 평가합니다.",
         "size": 10.5, "color": "1E293B"},
    ])

    # ── 카드 공통 치수 ─────────────────────────────────────────
    CY   = 1.97   # 카드 상단 y
    CH   = 2.47   # 카드 전체 높이
    HDRH = 0.42   # 컬러 헤더 높이
    BY   = CY + HDRH  # 바디 시작 y = 2.39

    # ── 카드 1: 지지도 (Support) ──────────────────────────────
    C1X, C1W = 0.45, 2.85
    rect(slide, C1X, CY, C1W, CH, "F0F7FF")
    rect_text(slide, C1X, CY, C1W, HDRH, "065A82",
              [{"text": "지지도  (Support)", "size": 12, "bold": True}],
              center=True)
    # 공식
    label(slide, C1X + 0.15, BY + 0.10, C1W - 0.30, 0.45, [
        {"text": "P(X ∩ Y)", "size": 20, "bold": True, "color": "065A82"},
    ])
    # 구분선
    rect(slide, C1X + 0.15, BY + 0.58, C1W - 0.30, 0.02, "BFDBFE")
    # 설명
    label(slide, C1X + 0.15, BY + 0.65, C1W - 0.30, 1.65, [
        {"text": "전체 PDF 중 두 키워드 X, Y가", "size": 9.5, "color": "374151"},
        {"text": "함께 등장하는 비율", "size": 9.5, "color": "374151"},
        {"text": "", "size": 5},
        {"text": "0에 가까울수록 희귀한 조합,", "size": 9, "color": "64748B"},
        {"text": "1에 가까울수록 자주 함께 등장", "size": 9, "color": "64748B"},
        {"text": "", "size": 5},
        {"text": "연관 규칙 탐색의 '필요 조건'", "size": 9, "bold": True, "color": "065A82"},
    ])

    # ── 카드 2: 신뢰도 (Confidence) ───────────────────────────
    C2X, C2W = 3.42, 2.85
    rect(slide, C2X, CY, C2W, CH, "FAF8FF")
    rect_text(slide, C2X, CY, C2W, HDRH, "7C3AED",
              [{"text": "신뢰도  (Confidence)", "size": 12, "bold": True}],
              center=True)
    label(slide, C2X + 0.15, BY + 0.08, C2W - 0.30, 0.60, [
        {"text": "P(Y|X)", "size": 20, "bold": True, "color": "7C3AED"},
        {"text": "= Support / P(X)", "size": 11, "color": "7C3AED"},
    ])
    rect(slide, C2X + 0.15, BY + 0.70, C2W - 0.30, 0.02, "DDD6FE")
    label(slide, C2X + 0.15, BY + 0.77, C2W - 0.30, 1.52, [
        {"text": "X 포함 PDF 중 Y도 함께 포함하는", "size": 9.5, "color": "374151"},
        {"text": "비율 — 연관 규칙의 정확도", "size": 9.5, "color": "374151"},
        {"text": "", "size": 5},
        {"text": "같은 지지도라면 신뢰도가 높은", "size": 9, "color": "64748B"},
        {"text": "규칙이 더 정확한 규칙", "size": 9, "color": "64748B"},
        {"text": "", "size": 5},
        {"text": "연관 규칙 탐색의 '충분 조건'", "size": 9, "bold": True, "color": "7C3AED"},
    ])

    # ── 카드 3: 향상도 (Lift) ─────────────────────────────────
    C3X, C3W = 6.39, 3.16
    rect(slide, C3X, CY, C3W, CH, "F0FDF9")
    rect_text(slide, C3X, CY, C3W, HDRH, "0D9488",
              [{"text": "향상도  (Lift)", "size": 12, "bold": True}],
              center=True)
    label(slide, C3X + 0.15, BY + 0.10, C3W - 0.30, 0.42, [
        {"text": "Confidence / P(Y)", "size": 15, "bold": True, "color": "0D9488"},
    ])
    rect(slide, C3X + 0.15, BY + 0.53, C3W - 0.30, 0.02, "A7F3D0")

    # 해석 배지 3개
    BX  = C3X + 0.15
    BW  = C3W - 0.30
    BH  = 0.43
    GAP = 0.12

    BY1 = BY + 0.60
    BY2 = BY1 + BH + GAP
    BY3 = BY2 + BH + GAP

    rect_text(slide, BX, BY1, BW, BH, "064E3B",   # 짙은 초록
              [{"text": "Lift  >  1  →  양의 상관관계 — 연관성 있음  ✓",
                "size": 9.5, "bold": True}])
    rect_text(slide, BX, BY2, BW, BH, "374151",   # 슬레이트
              [{"text": "Lift  =  1  →  독립 관계 — 연관 없음",
                "size": 9.5, "bold": False}])
    rect_text(slide, BX, BY3, BW, BH, "7C2D12",   # 짙은 오렌지
              [{"text": "Lift  <  1  →  음의 상관관계 — 역관계",
                "size": 9.5, "bold": False}])

    # ── 하단 다크 스트립 (예시 + 적용) ───────────────────────
    SY = CY + CH + 0.12   # = 4.56
    rect(slide, 0.45, SY, 9.10, 0.65, "1B2A4A")
    rect(slide, 0.45, SY, 0.05, 0.65, "0D9488")
    label(slide, 0.65, SY + 0.05, 8.80, 0.55, [
        {"text": "예시:  {pipeline, docker}  →  Support=0.67 · Confidence=0.80 · Lift=1.5",
         "size": 10, "color": "D1D5DB"},
        {"text": "Lift > 1  →  양의 상관관계  →  연관 규칙 점수(×2.5) + Jaccard 유사도 합산  →  상위 5개 PDF 추천",
         "size": 10, "bold": True, "color": "6EE7B7"},
    ])


# ── 실행 ─────────────────────────────────────────────────────────────────

def main():
    prs = Presentation(str(PPTX_PATH))
    print(f"기존 슬라이드 수: {len(prs.slides)}")

    target = prs.slides[6]  # 인덱스 6 = 슬라이드 7 (Apriori 슬라이드)
    clear_slide(target)
    build(target)

    out_path = PPTX_PATH.with_stem(PPTX_PATH.stem + "_updated")
    prs.save(str(out_path))
    print(f"완료: {out_path.name}  (총 {len(prs.slides)}장)")
    print("원본 파일이 열려있으면 닫은 뒤 파일명을 변경하거나, _updated 버전을 사용하세요.")


if __name__ == "__main__":
    main()
