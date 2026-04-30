"""Create system architecture slide for Mail2Task presentation."""
import sys, io as sysio
sys.stdout = sysio.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import requests, os, tempfile, copy
from io import BytesIO
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO
from pptx.oxml.ns import qn


def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


tmp = tempfile.gettempdir()

# Download icons
hdrs = {'User-Agent': 'Mozilla/5.0'}
ICON_URLS = {
    'gmail':   'https://img.icons8.com/color/96/gmail-new.png',
    'python':  'https://img.icons8.com/color/96/python.png',
    'claude':  'https://img.icons8.com/color/96/artificial-intelligence.png',
    'mongodb': 'https://img.icons8.com/color/96/mongodb.png',
    'webui':   'https://img.icons8.com/color/96/domain.png',
}

icon_paths = {}
for name, url in ICON_URLS.items():
    path = os.path.join(tmp, f'arch_{name}.png')
    if not os.path.exists(path):
        r = requests.get(url, timeout=10, headers=hdrs)
        img = Image.open(BytesIO(r.content)).resize((72, 72), Image.LANCZOS)
        img.save(path, 'PNG')
    print(f'Icon ready: {name}')
    icon_paths[name] = path

# Open presentation
prs = Presentation('Mail2Task_Presentation_v2_recommendation.pptx')
print(f'Original slide count: {len(prs.slides)}')

# Find blank layout (this presentation has only one layout)
blank_layout = prs.slide_layouts[0]

# Add new slide and move to index 6 (after 04 system structure at index 5)
new_slide = prs.slides.add_slide(blank_layout)
sldIdLst = prs.part._element.find(qn('p:sldIdLst'))
slide_id_elems = list(sldIdLst)
new_sld_id = slide_id_elems[-1]
sldIdLst.remove(new_sld_id)
sldIdLst.insert(6, new_sld_id)
new_slide = prs.slides[6]
print(f'New slide at index 6. Total: {len(prs.slides)}')

# Copy header (first 9 shapes) from slide index 5 (04 system structure)
source_slide = prs.slides[5]
source_shapes = list(source_slide.shapes)
for i in range(9):
    sp_copy = copy.deepcopy(source_shapes[i]._element)
    new_slide.shapes._spTree.append(sp_copy)
print('Header copied')

# Update title: runs[3] = " 구조" -> " 아키텍처"
for shape in new_slide.shapes:
    if shape.name == 'Text 3' and shape.has_text_frame:
        runs = list(shape.text_frame.paragraphs[0].runs)
        if len(runs) >= 4:
            runs[3].text = ' 아키텍처'
        break

# Update subtitle (Text 6)
for shape in new_slide.shapes:
    if shape.name == 'Text 6' and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            rl = list(para.runs)
            if rl:
                rl[0].text = 'Gmail IMAP · FastAPI · OpenAI API · MongoDB 핵심 컴포넌트 구성도'
                for run in rl[1:]:
                    run.text = ''
        break

# Update page number (Text 8)
for shape in new_slide.shapes:
    if shape.name == 'Text 8' and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = '7  /  14'
        break

print('Header text updated')

# ── Layout constants ───────────────────────────────────────────────────────────
BOX_W = 2.1
BOX_H = 1.42
C_GMAIL = 0.35
C_PIPE  = 3.75
C_AI    = 7.15
R1 = 1.58
R2 = 3.60


def add_box(slide, left, top, color, icon_path, title, desc):
    # Rounded rectangle background
    bg = slide.shapes.add_shape(
        MSO.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(BOX_W), Inches(BOX_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(color)
    bg.line.fill.background()

    # White circle behind icon
    ICON_SZ = 0.56
    ix = left + (BOX_W - ICON_SZ) / 2
    iy = top + 0.1
    circ = slide.shapes.add_shape(
        MSO.OVAL,
        Inches(ix - 0.06), Inches(iy - 0.06),
        Inches(ICON_SZ + 0.12), Inches(ICON_SZ + 0.12))
    circ.fill.solid()
    circ.fill.fore_color.rgb = rgb('FFFFFF')
    circ.line.fill.background()

    # Icon image
    slide.shapes.add_picture(
        icon_path, Inches(ix), Inches(iy), Inches(ICON_SZ), Inches(ICON_SZ))

    # Title
    tb1 = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.76), Inches(BOX_W), Inches(0.33))
    tb1.text_frame.word_wrap = False
    p1 = tb1.text_frame.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = rgb('FFFFFF')

    # Description
    tb2 = slide.shapes.add_textbox(
        Inches(left + 0.08), Inches(top + 1.1), Inches(BOX_W - 0.16), Inches(0.28))
    tb2.text_frame.word_wrap = True
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(7.5)
    r2.font.bold = False
    r2.font.color.rgb = rgb('C8E4FA')


def add_arrow_h(slide, x1, x2, y):
    w, h = x2 - x1, 0.2
    arr = slide.shapes.add_shape(
        MSO.RIGHT_ARROW,
        Inches(x1), Inches(y - h / 2), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = rgb('0D9488')
    arr.line.fill.background()


def add_arrow_v(slide, x, y1, y2):
    h, w = y2 - y1, 0.2
    arr = slide.shapes.add_shape(
        MSO.DOWN_ARROW,
        Inches(x - w / 2), Inches(y1), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = rgb('0D9488')
    arr.line.fill.background()


def add_label(slide, text, left, top):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(BOX_W), Inches(0.22))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(7.5)
    r.font.color.rgb = rgb('94A3B8')


# ── Layer labels ───────────────────────────────────────────────────────────────
add_label(new_slide, '[ 입력 ]',       C_GMAIL, R1 - 0.27)
add_label(new_slide, '[ 처리 엔진 ]',  C_PIPE,  R1 - 0.27)
add_label(new_slide, '[ AI 서비스 ]',  C_AI,    R1 - 0.27)
add_label(new_slide, '[ 데이터 저장 ]', C_PIPE, R2 - 0.27)
add_label(new_slide, '[ UI 출력 ]',    C_AI,    R2 - 0.27)

# ── Component boxes ────────────────────────────────────────────────────────────
add_box(new_slide, C_GMAIL, R1, 'B91C1C', icon_paths['gmail'],
        'Gmail IMAP', 'imaplib · 메일 수신 · PDF 첨부')
add_box(new_slide, C_PIPE,  R1, '1E3A5F', icon_paths['python'],
        'Core Pipeline', 'mail/ · core/ · tasks/ 모듈')
add_box(new_slide, C_AI,    R1, '6D28D9', icon_paths['claude'],
        'OpenAI API', '작업 요약 자동 생성 (summarizer)')
add_box(new_slide, C_PIPE,  R2, '166534', icon_paths['mongodb'],
        'MongoDB', 'mails · tasks · pdf_documents')
add_box(new_slide, C_AI,    R2, '0C4A6E', icon_paths['webui'],
        'FastAPI Web UI', 'Jinja2 · 대시보드 · 작업 관리')
print('Component boxes added')

# ── Arrows ─────────────────────────────────────────────────────────────────────
ROW1_MID = R1 + BOX_H / 2
ROW2_MID = R2 + BOX_H / 2
PIPE_CX  = C_PIPE + BOX_W / 2

add_arrow_h(new_slide, C_GMAIL + BOX_W, C_PIPE,      ROW1_MID)  # Gmail → Pipeline
add_arrow_h(new_slide, C_PIPE  + BOX_W, C_AI,        ROW1_MID)  # Pipeline → OpenAI
add_arrow_v(new_slide, PIPE_CX,          R1 + BOX_H, R2)         # Pipeline ↓ MongoDB
add_arrow_h(new_slide, C_PIPE  + BOX_W, C_AI,        ROW2_MID)  # MongoDB → Web UI
print('Arrows added')

# Save
output = 'Mail2Task_Presentation_v2_recommendation.pptx'
prs.save(output)
print(f'\nSaved: {output}')

# Verify
prs2 = Presentation(output)
print(f'Final slide count: {len(prs2.slides)}')
